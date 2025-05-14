import PyPDF2
from docx import Document
from pptx import Presentation
import pdfminer.high_level
from typing import Optional
from app.core.constants import ALLOWED_FILE_TYPES

class FileProcessor:
    @staticmethod
    def process_file(file_path: str) -> Optional[str]:
        extension = file_path.split('.')[-1].lower()
        
        if extension not in ALLOWED_FILE_TYPES:
            raise ValueError(f"Unsupported file type: {extension}")

        try:
            if extension == 'pdf':
                return FileProcessor._read_pdf(file_path)
            elif extension == 'docx':
                return FileProcessor._read_docx(file_path)
            elif extension == 'txt':
                return FileProcessor._read_text(file_path)
            elif extension == 'pptx':
                return FileProcessor._read_pptx(file_path)
        except Exception as e:
            raise RuntimeError(f"Error processing file: {str(e)}")

    @staticmethod
    def _read_pdf(file_path: str) -> str:
        return pdfminer.high_level.extract_text(file_path)

    @staticmethod
    def _read_docx(file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    @staticmethod
    def _read_text(file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    @staticmethod
    def _read_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)